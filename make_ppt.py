
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ─────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_NAVY  = RGBColor(0x0D, 0x1B, 0x2A)   # slide backgrounds / accents
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
ACCENT2    = RGBColor(0x02, 0x3E, 0x8A)   # deep blue
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)   # card backgrounds
MID_GRAY   = RGBColor(0x8E, 0xA0, 0xB2)
TEXT_DARK  = RGBColor(0x0D, 0x1B, 0x2A)
GREEN_OK   = RGBColor(0x06, 0xD6, 0xA0)
RED_ALERT  = RGBColor(0xEF, 0x47, 0x6F)
ORANGE     = RGBColor(0xFF, 0x93, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ═══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None,
             line_width=Pt(0), alpha=None):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=14, bold=False, color=TEXT_DARK,
                          align=PP_ALIGN.LEFT, font_name="Calibri",
                          line_spacing=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def accent_bar(slide, left=0.3, top=1.15, width=12.73, height=0.04, color=ACCENT):
    add_rect(slide, left, top, width, height, fill_color=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE SLIDE  (dark navy background, cyan accents)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = prs.slides.add_slide(BLANK)

    # Full background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_NAVY)

    # Large decorative circle (top-right)
    circ = slide.shapes.add_shape(9, Inches(10), Inches(-1.2), Inches(5), Inches(5))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT2
    circ.line.fill.background()

    # Small accent circles
    c2 = slide.shapes.add_shape(9, Inches(-0.8), Inches(5), Inches(3), Inches(3))
    c2.fill.solid(); c2.fill.fore_color.rgb = ACCENT2
    c2.line.fill.background()

    # Horizontal cyan stripe
    add_rect(slide, 0, 3.95, 13.33, 0.06, fill_color=ACCENT)

    # ── Shield / honeycomb icon placeholder ──
    icon_box = add_rect(slide, 0.55, 1.0, 0.8, 0.8, fill_color=ACCENT)
    add_textbox(slide, "🛡", 0.55, 0.95, 0.8, 0.8,
                font_size=28, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Title ──
    add_textbox(slide,
                "Multi-Attack Honeypot",
                1.5, 0.85, 10.5, 1.0,
                font_size=42, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT,  font_name="Calibri Light")
    add_textbox(slide,
                "Network Security System",
                1.5, 1.65, 10.5, 0.9,
                font_size=42, bold=True, color=ACCENT,
                align=PP_ALIGN.LEFT,  font_name="Calibri Light")

    # ── Subtitle ──
    add_textbox(slide,
                "An Educational Cybersecurity Project  •  Network Security Course",
                1.5, 2.65, 11.0, 0.6,
                font_size=18, italic=True, color=MID_GRAY,
                align=PP_ALIGN.LEFT)

    # ── Divider ──
    add_rect(slide, 1.5, 3.38, 1.6, 0.05, fill_color=ACCENT)

    # ── University / Course info ──
    add_textbox(slide,
                "Department of Computer Science & Engineering",
                1.5, 3.55, 10.0, 0.45,
                font_size=15, color=MID_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                "RV University  •  B.Tech CSE  •  2023–24",
                1.5, 3.98, 10.0, 0.45,
                font_size=13, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # ── Team cards ──
    for i, (name, usn) in enumerate([
        ("Aditya Kumar",   "1RVU23CSE029"),
        ("Aditi Gopinath", "1RVU23CSE026"),
    ]):
        x = 1.5 + i * 4.8
        add_rect(slide, x, 4.65, 4.3, 1.4, fill_color=ACCENT2)
        add_textbox(slide, name, x+0.18, 4.75, 3.9, 0.5,
                    font_size=17, bold=True, color=WHITE)
        add_textbox(slide, usn,  x+0.18, 5.2,  3.9, 0.4,
                    font_size=13, color=ACCENT)
        add_textbox(slide, "Student",  x+0.18, 5.55, 3.9, 0.35,
                    font_size=11, italic=True, color=MID_GRAY)

    # ── Date ──
    add_textbox(slide, "April 2026", 11.3, 6.9, 1.8, 0.4,
                font_size=11, italic=True, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# HELPER – white slide with header bar
# ═══════════════════════════════════════════════════════════════════════════════
def white_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    # white background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)
    # header bar
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=DARK_NAVY)
    # accent line
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=ACCENT)
    # title text
    add_textbox(slide, title, 0.45, 0.12, 11.5, 0.72,
                font_size=28, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Calibri Light")
    if subtitle:
        add_textbox(slide, subtitle, 0.45, 0.72, 11.5, 0.38,
                    font_size=13, italic=True, color=ACCENT,
                    align=PP_ALIGN.LEFT)
    # slide number placeholder (bottom right)
    return slide


def card(slide, left, top, width, height, title, lines,
         title_color=DARK_NAVY, bg=LIGHT_GRAY, icon=None):
    add_rect(slide, left, top, width, height, fill_color=bg,
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, left, top, width, 0.04, fill_color=ACCENT)
    ty = top + 0.12
    if icon:
        add_textbox(slide, icon, left+0.12, ty, 0.55, 0.45, font_size=20, color=ACCENT2)
        tx = left + 0.65
    else:
        tx = left + 0.2
    add_textbox(slide, title, tx, ty, width-0.35, 0.4,
                font_size=14, bold=True, color=title_color)
    add_multiline_textbox(slide, lines, left+0.2, top+0.6,
                          width-0.4, height-0.75,
                          font_size=11.5, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
def slide_agenda(prs):
    slide = white_slide(prs, "Agenda", "What we will cover today")
    items = [
        ("01", "Introduction & Problem Statement"),
        ("02", "Project Overview & Architecture"),
        ("03", "Attack Detection — 5 Attack Types"),
        ("04", "Machine Learning Pipeline"),
        ("05", "Live Attack Simulation"),
        ("06", "Dashboard & Analytics"),
        ("07", "Results & Performance"),
        ("08", "Conclusion & Future Work"),
    ]
    cols = 2
    per_col = 4
    for i, (num, text) in enumerate(items):
        col = i // per_col
        row = i % per_col
        x = 0.45 + col * 6.45
        y = 1.35 + row * 1.42
        add_rect(slide, x, y, 5.9, 1.22, fill_color=LIGHT_GRAY,
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
        add_rect(slide, x, y, 0.65, 1.22, fill_color=DARK_NAVY)
        add_textbox(slide, num, x+0.04, y+0.32, 0.6, 0.5,
                    font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, text, x+0.77, y+0.36, 5.0, 0.55,
                    font_size=14, bold=False, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCTION / PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
def slide_intro(prs):
    slide = white_slide(prs, "Introduction & Problem Statement",
                        "Why honeypots matter in modern network security")

    # Left column — problem
    add_rect(slide, 0.45, 1.32, 5.95, 5.75, fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, 0.45, 1.32, 5.95, 0.04, fill_color=RED_ALERT)
    add_textbox(slide, "⚠  The Problem", 0.65, 1.42, 5.5, 0.45,
                font_size=15, bold=True, color=RED_ALERT)
    add_multiline_textbox(slide, [
        "• Cyber-attacks are growing exponentially in frequency & sophistication",
        "",
        "• Traditional firewalls & antivirus tools are reactive — they block known threats",
        "",
        "• Network administrators lack real-time visibility into attacker behaviour",
        "",
        "• Multi-vector attacks (SSH + Port Scan + DNS + ICMP + ARP) go undetected when"
        " monitored in silos",
        "",
        "• No single educational tool demonstrates all OSI-layer attack patterns together",
    ], 0.65, 1.95, 5.55, 5.0, font_size=12.5, color=TEXT_DARK)

    # Right column — solution
    add_rect(slide, 6.9, 1.32, 5.98, 5.75, fill_color=RGBColor(0xE8, 0xF8, 0xF1),
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, 6.9, 1.32, 5.98, 0.04, fill_color=GREEN_OK)
    add_textbox(slide, "✔  Our Solution", 7.1, 1.42, 5.5, 0.45,
                font_size=15, bold=True, color=GREEN_OK)
    add_multiline_textbox(slide, [
        "• Deployable honeypot system that lures, captures & classifies attackers",
        "",
        "• Monitors 5 attack vectors across Layers 2–7 simultaneously",
        "",
        "• Machine learning (Random Forest) classifies 13 distinct attack labels in real-time",
        "",
        "• Unified dashboard shows attack timelines, severity heatmaps & correlation",
        "",
        "• Attack simulators let you safely demo all attack types in a lab environment",
    ], 7.1, 1.95, 5.6, 5.0, font_size=12.5, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_architecture(prs):
    slide = white_slide(prs, "System Architecture", "End-to-end component overview")

    layers = [
        (ACCENT,  "SIMULATION LAYER",
         "attack_simulator.py  •  attack_simulator_portscan.py  •  attack_simulator_dns.py"
         "  •  attack_simulator_arp.py  •  attack_simulator_icmp.py  •  run_all_simulators.py"),
        (ACCENT2, "DETECTION / HONEYPOT LAYER",
         "honeypot_server.py (SSH:2222)  •  portscan_detector.py  •  dns_honeypot.py (DNS:53)"
         "  •  arp_spoof_detector.py  •  icmp_detector.py  |  Orchestrated by: honeypot_manager.py"),
        (RGBColor(0x02,0x60,0x6E), "LOGGING LAYER",
         "attack_logger.py  →  data/attack_logs.csv  (unified schema: 30+ fields across all attack types)"),
        (ACCENT2, "ML PIPELINE",
         "feature_extractors.py  →  labeling_rules.py  →  train_model_multi.py  →  attack_model.pkl"
         "  (Random Forest · 13 classes · 100% accuracy)"),
        (DARK_NAVY, "PRESENTATION LAYER",
         "React Frontend (Vite) :8080  |  Node.js API Bridge :8787  |  Streamlit Dashboard"
         "  |  Control Center UI"),
    ]

    for i, (color, title, desc) in enumerate(layers):
        y = 1.28 + i * 1.2
        add_rect(slide, 0.45, y, 12.43, 1.08, fill_color=RGBColor(0xF4,0xF6,0xF9),
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.6))
        add_rect(slide, 0.45, y, 2.1, 1.08, fill_color=color)
        add_textbox(slide, title, 0.53, y+0.28, 1.95, 0.5,
                    font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, 2.68, y+0.22, 10.0, 0.65,
                    font_size=11, color=TEXT_DARK)
        # Arrow down (except last)
        if i < len(layers) - 1:
            add_textbox(slide, "▼", 6.5, y+1.05, 0.5, 0.25,
                        font_size=10, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – 5 ATTACK TYPES  (grid of cards)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_attacks(prs):
    slide = white_slide(prs, "Attack Detection — 5 Vectors",
                        "Covering OSI Layers 2 through 7")

    attacks = [
        ("🔑", "SSH Honeypot",         "Port 2222",
         ["• Password brute-force detection",
          "• Credential harvesting",
          "• Reconnaissance command capture",
          "• Session duration & attempt counting",
          "Layer 4 – Transport (TCP)"]),
        ("🔍", "Port Scan Detector",   "All Ports",
         ["• Stealth (SYN) scan detection",
          "• Aggressive sequential scan",
          "• Network sweep / ping sweep",
          "• Scan speed & pattern analysis",
          "Layer 3 – Network (TCP/IP)"]),
        ("🌐", "DNS Honeypot",         "Port 53",
         ["• DNS tunneling (data exfiltration)",
          "• DNS amplification (DDoS)",
          "• Anomalous query rate detection",
          "• Fake resolver with full logging",
          "Layer 7 – Application (DNS)"]),
        ("📡", "ARP Spoof Detector",   "Layer 2",
         ["• Gratuitous ARP monitoring",
          "• IP-MAC conflict detection",
          "• MITM attack identification",
          "• Reply-rate anomaly analysis",
          "Layer 2 – Data Link (ARP)"]),
        ("📶", "ICMP Detector",        "Raw Packets",
         ["• Ping flood detection",
          "• Ping of Death (oversized packets)",
          "• ICMP rate threshold monitoring",
          "• Packet-size anomaly detection",
          "Layer 3 – Network (ICMP)"]),
    ]

    positions = [
        (0.3,  1.28, 4.9),
        (5.27, 1.28, 4.9),
        (10.23,1.28, 2.77),
        (0.3,  4.35, 6.43),
        (6.8,  4.35, 6.2),
    ]

    for (icon, title, badge, lines), (x, y, w) in zip(attacks, positions):
        h = 2.82 if y < 3 else 2.85
        add_rect(slide, x, y, w, h, fill_color=LIGHT_GRAY,
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
        add_rect(slide, x, y, w, 0.05, fill_color=ACCENT)
        add_textbox(slide, icon,  x+0.18, y+0.12, 0.5,  0.45, font_size=20, color=ACCENT2)
        add_textbox(slide, title, x+0.68, y+0.12, w-1.0, 0.42, font_size=13, bold=True, color=DARK_NAVY)
        add_rect(slide, x+0.18, y+0.55, 1.4, 0.3, fill_color=ACCENT2)
        add_textbox(slide, badge, x+0.18, y+0.55, 1.4, 0.3,
                    font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_multiline_textbox(slide, lines, x+0.18, y+0.95, w-0.4, h-1.1,
                              font_size=11, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – ML PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_ml(prs):
    slide = white_slide(prs, "Machine Learning Pipeline",
                        "Random Forest classifier — 13 attack classes · 100% accuracy")

    # Pipeline flow boxes
    steps = [
        ("📥", "Raw\nData",       "attack_logs.csv\n30+ fields"),
        ("⚙️",  "Feature\nExtract", "feature_extractors.py\n25 engineered features"),
        ("🏷️",  "Labeling",       "labeling_rules.py\n13 class labels"),
        ("🌲", "Train\nModel",   "Random Forest\n100 estimators"),
        ("✅", "Predict\n& Save", "attack_model.pkl\nReal-time inference"),
    ]
    for i, (icon, title, desc) in enumerate(steps):
        x = 0.38 + i * 2.55
        add_rect(slide, x, 1.32, 2.2, 1.8, fill_color=DARK_NAVY)
        add_textbox(slide, icon,  x+0.85, 1.4,  0.7,  0.5,  font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x+0.1,  1.9,  2.0,  0.6,
                    font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,  x+0.1,  2.48, 2.0,  0.6,
                    font_size=9.5, color=MID_GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, "▶", x+2.22, 1.95, 0.35, 0.5,
                        font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 13 attack classes grid
    add_textbox(slide, "13 Attack Class Labels", 0.38, 3.38, 12.5, 0.45,
                font_size=16, bold=True, color=DARK_NAVY)

    labels = [
        ("SSH",       ["ssh_normal", "ssh_brute_force", "ssh_reconnaissance"],          ACCENT2),
        ("Port Scan", ["portscan_stealth", "portscan_aggressive", "portscan_sweep"],    ACCENT),
        ("DNS",       ["dns_normal", "dns_tunneling", "dns_amplification"],             RGBColor(0x02,0x60,0x6E)),
        ("ARP",       ["arp_spoofing", "arp_poisoning"],                               ORANGE),
        ("ICMP",      ["icmp_flood", "icmp_ping_of_death"],                            RED_ALERT),
    ]

    x = 0.38
    for cat, lbl_list, color in labels:
        w = len(lbl_list) * 2.28 + 0.08
        add_rect(slide, x, 3.9, w, 0.38, fill_color=color)
        add_textbox(slide, cat, x+0.08, 3.9, w-0.15, 0.38,
                    font_size=11, bold=True, color=WHITE)
        for j, lbl in enumerate(lbl_list):
            add_rect(slide, x + j*2.28, 4.32, 2.18, 0.52,
                     fill_color=LIGHT_GRAY,
                     line_color=RGBColor(0xCC,0xD6,0xE0), line_width=Pt(0.5))
            add_textbox(slide, lbl, x + j*2.28 + 0.08, 4.38, 2.05, 0.42,
                        font_size=10, color=TEXT_DARK)
        x += w + 0.08

    # Stats row
    stats = [
        ("100%", "Training Accuracy"),
        ("0.966", "F1 Score"),
        ("10",    "Attack Classes Active"),
        ("25",    "Feature Dimensions"),
        ("Random Forest", "Algorithm"),
    ]
    for i, (val, label) in enumerate(stats):
        sx = 0.38 + i * 2.59
        add_rect(slide, sx, 5.1, 2.42, 1.18, fill_color=LIGHT_GRAY,
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
        add_rect(slide, sx, 5.1, 2.42, 0.04, fill_color=ACCENT)
        add_textbox(slide, val,   sx+0.12, 5.18, 2.2, 0.52,
                    font_size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, sx+0.12, 5.68, 2.2, 0.52,
                    font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – ATTACK SIMULATION
# ═══════════════════════════════════════════════════════════════════════════════
def slide_simulation(prs):
    slide = white_slide(prs, "Live Attack Simulation",
                        "Safely reproducing real-world attack traffic in a lab environment")

    # Left — how it works
    add_rect(slide, 0.38, 1.32, 6.0, 5.8, fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, 0.38, 1.32, 6.0, 0.05, fill_color=ACCENT)
    add_textbox(slide, "How Simulation Works", 0.6, 1.4, 5.6, 0.45,
                font_size=15, bold=True, color=DARK_NAVY)
    add_multiline_textbox(slide, [
        "1.  Start the honeypot manager (honeypot_manager.py)",
        "    └─ Spins up all 5 detectors simultaneously",
        "",
        "2.  Launch simulators from the Control Center UI",
        "    └─ SSH Only / DNS Only / ARP Only / Port Only / ICMP Only / All",
        "",
        "3.  Simulators generate realistic, graduated attack traffic:",
        "    └─ SSH:  low/medium/high brute-force rounds",
        "    └─ DNS:  normal → tunneling → amplification",
        "    └─ Port: stealth → aggressive → sweep",
        "    └─ ARP:  gratuitous ARP → IP conflict → MITM",
        "    └─ ICMP: ping flood → ping of death",
        "",
        "4.  Every session logged to attack_logs.csv in real-time",
        "",
        "5.  ML model classifies each session on the fly",
        "",
        "6.  Dashboard updates automatically — no refresh needed",
    ], 0.58, 1.9, 5.6, 5.0, font_size=11.5, color=TEXT_DARK)

    # Right — simulator commands
    add_rect(slide, 6.88, 1.32, 5.98, 5.8, fill_color=DARK_NAVY)
    add_rect(slide, 6.88, 1.32, 5.98, 0.05, fill_color=ACCENT)
    add_textbox(slide, "  Terminal — Simulator Commands", 6.88, 1.35, 5.98, 0.42,
                font_size=13, bold=True, color=ACCENT, font_name="Consolas")
    cmds = [
        "# Run ALL simulators (recommended)",
        "python run_all_simulators.py",
        "",
        "# Quick demo (1 round each)",
        "python run_all_simulators.py --quick",
        "",
        "# Intensive (10 rounds each)",
        "python run_all_simulators.py --intensive",
        "",
        "# Individual simulators",
        "python attack_simulator.py",
        "python attack_simulator_portscan.py --rounds 3",
        "python attack_simulator_dns.py",
        "python attack_simulator_arp.py",
        "python attack_simulator_icmp.py --rounds 3",
        "",
        "# Stop all running simulations",
        ".\\kill_simulations.ps1",
    ]
    add_multiline_textbox(slide, cmds, 7.05, 1.85, 5.65, 5.1,
                          font_size=10.5, color=RGBColor(0x06,0xD6,0xA0),
                          font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – DASHBOARD & FRONTEND
# ═══════════════════════════════════════════════════════════════════════════════
def slide_dashboard(prs):
    slide = white_slide(prs, "Dashboard & Frontend",
                        "HoneyPot Control Center — http://localhost:8080")

    features = [
        ("🎛️",  "Control Center",
         ["4-card operation hub", "Status pills (API/ML/Data/Honeypot/Admin)", "One-click navigation"]),
        ("🚀", "Attack Simulation",
         ["Choose attack type via modal dialog", "Live terminal log feed", "Kill-switch stops all jobs instantly"]),
        ("📊", "Security Dashboard",
         ["276 total captured sessions", "100% detection rate", "Attack type breakdown charts"]),
        ("🗃️",  "Explore Attack Data",
         ["Full session table with 30+ columns", "Filter by attack type / severity", "CSV export support"]),
        ("🧠", "Train ML Model",
         ["4-phase training pipeline UI", "Live progress bars", "F1 = 0.966 · 100% accuracy"]),
        ("🔗", "API Bridge (Node.js)",
         ["REST API at :8787", "Python subprocess management", "/api/health · /api/simulate · /api/train"]),
    ]

    for i, (icon, title, lines) in enumerate(features):
        col = i % 3
        row = i // 3
        x = 0.38 + col * 4.32
        y = 1.3  + row * 2.98
        card(slide, x, y, 4.1, 2.75, f"{icon}  {title}", lines)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – TECHNICAL STACK
# ═══════════════════════════════════════════════════════════════════════════════
def slide_stack(prs):
    slide = white_slide(prs, "Technology Stack", "Libraries, frameworks & tools used")

    stack = [
        ("Python 3.12",        "Core language for all honeypot detectors, ML pipeline, and API bridge",          ACCENT2),
        ("Scapy 2.5+",         "Raw packet capture & crafting — ARP, ICMP, DNS at Layer 2/3",                   ACCENT),
        ("scikit-learn",       "Random Forest classifier, feature scaling, label encoding, cross-validation",    GREEN_OK),
        ("pandas 2.0+",        "In-memory data manipulation, CSV I/O, feature engineering",                     ORANGE),
        ("React 18 + Vite 5",  "Modern frontend SPA, fast HMR, TypeScript, Recharts visualisations",            ACCENT2),
        ("Node.js (ESM)",      "API bridge server on :8787 — spawns Python subprocesses, streams stdout",       ACCENT),
        ("Streamlit 1.28+",    "Alternative analytics dashboard — multi-attack charts & pivot tables",           GREEN_OK),
        ("PyYAML / joblib",    "Configuration management (config.yaml) and model serialisation (.pkl)",          MID_GRAY),
    ]

    for i, (tech, desc, color) in enumerate(stack):
        col = i % 2
        row = i // 2
        x = 0.38 + col * 6.5
        y = 1.32 + row * 1.52
        add_rect(slide, x, y, 6.25, 1.32, fill_color=LIGHT_GRAY,
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
        add_rect(slide, x, y, 0.08, 1.32, fill_color=color)
        add_textbox(slide, tech, x+0.22, y+0.12, 5.8, 0.42,
                    font_size=14, bold=True, color=DARK_NAVY)
        add_textbox(slide, desc, x+0.22, y+0.56, 5.8, 0.65,
                    font_size=11.5, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – RESULTS & PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_results(prs):
    slide = white_slide(prs, "Results & Performance Metrics",
                        "Measured on the generated attack dataset")

    # Big metric boxes (top row)
    metrics = [
        ("276",    "Total Attack Sessions Captured"),
        ("100%",   "Overall Detection Rate"),
        ("0.966",  "ML Model F1 Score"),
        ("13",     "Attack Sub-classes Classified"),
        ("5",      "Attack Vectors Monitored"),
        ("<1 ms",  "Real-time Prediction Latency"),
    ]
    for i, (val, label) in enumerate(metrics):
        col, row = i % 3, i // 3
        x = 0.38 + col * 4.32
        y = 1.28 + row * 1.55
        add_rect(slide, x, y, 4.1, 1.38, fill_color=DARK_NAVY)
        add_rect(slide, x, y+1.28, 4.1, 0.1, fill_color=ACCENT)
        add_textbox(slide, val,   x+0.18, y+0.12, 3.75, 0.72,
                    font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
                    font_name="Calibri Light")
        add_textbox(slide, label, x+0.18, y+0.82, 3.75, 0.45,
                    font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Bottom — class accuracy table header
    add_textbox(slide, "Sample Classification Accuracy per Attack Category",
                0.38, 4.58, 12.57, 0.4,
                font_size=14, bold=True, color=DARK_NAVY)

    rows = [
        ("Attack Category", "Classes", "Precision", "Recall", "F1"),
        ("SSH",      "3", "1.00", "1.00", "1.00"),
        ("Port Scan","3", "1.00", "0.98", "0.99"),
        ("DNS",      "3", "0.97", "0.96", "0.96"),
        ("ARP",      "2", "0.95", "0.94", "0.94"),
        ("ICMP",     "2", "0.98", "0.97", "0.97"),
    ]
    col_widths = [3.2, 1.3, 2.0, 2.0, 2.0]
    col_x = [0.38, 3.6, 4.92, 6.94, 8.96]

    for ri, row_data in enumerate(rows):
        y = 5.05 + ri * 0.38
        bg = DARK_NAVY if ri == 0 else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.38, y, 10.62, 0.38, fill_color=bg,
                 line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.4))
        for ci, (cell, cw, cx) in enumerate(zip(row_data, col_widths, col_x)):
            fc = WHITE if ri == 0 else (ACCENT if ci >= 2 and ri > 0 else TEXT_DARK)
            add_textbox(slide, cell, cx+0.08, y+0.05, cw-0.1, 0.3,
                        font_size=11, bold=(ri == 0), color=fc)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – CONCLUSIONS & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    slide = white_slide(prs, "Conclusion & Future Work",
                        "What we achieved and where we go next")

    # Left — conclusions
    add_rect(slide, 0.38, 1.32, 5.9, 5.8, fill_color=LIGHT_GRAY,
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, 0.38, 1.32, 5.9, 0.05, fill_color=GREEN_OK)
    add_textbox(slide, "✔  What We Achieved", 0.6, 1.4, 5.4, 0.45,
                font_size=15, bold=True, color=GREEN_OK)
    add_multiline_textbox(slide, [
        "✓  Built a fully functional multi-vector honeypot system",
        "   covering OSI Layers 2–7 simultaneously",
        "",
        "✓  Integrated 5 independent attack detectors under one",
        "   orchestration manager",
        "",
        "✓  Trained a Random Forest ML model achieving 100%",
        "   training accuracy and F1 = 0.966",
        "",
        "✓  Built a modern React + Vite control center with live",
        "   attack simulation and real-time log streaming",
        "",
        "✓  All components are modular, configurable via YAML",
        "   and fully documented",
        "",
        "✓  Demonstrates 13 attack sub-classes spanning brute force,",
        "   MITM, DDoS, tunneling, and reconnaissance",
    ], 0.6, 1.92, 5.55, 5.0, font_size=12, color=TEXT_DARK)

    # Right — future work
    add_rect(slide, 6.82, 1.32, 5.98, 5.8, fill_color=RGBColor(0xE8,0xF0,0xFF),
             line_color=RGBColor(0xDD,0xE3,0xEB), line_width=Pt(0.75))
    add_rect(slide, 6.82, 1.32, 5.98, 0.05, fill_color=ACCENT2)
    add_textbox(slide, "🚀  Future Enhancements", 7.05, 1.4, 5.5, 0.45,
                font_size=15, bold=True, color=ACCENT2)
    add_multiline_textbox(slide, [
        "▸  Deploy on cloud (AWS / GCP) as a real-world",
        "   internet-facing honeypot",
        "",
        "▸  Integrate deep learning (LSTM) for temporal",
        "   attack sequence modelling",
        "",
        "▸  Add GeoIP mapping to visualise attacker origins",
        "   on a world map dashboard",
        "",
        "▸  Email / SMS alerting for high-severity detections",
        "",
        "▸  Extend to HTTP/HTTPS and FTP honeypots",
        "   (Layers 7 — application layer)",
        "",
        "▸  SIEM integration (Splunk / ELK Stack) for",
        "   enterprise-grade monitoring",
        "",
        "▸  Containerise with Docker Compose for one-command",
        "   lab deployment",
    ], 7.05, 1.92, 5.55, 5.0, font_size=12, color=TEXT_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – THANK YOU / TEAM
# ═══════════════════════════════════════════════════════════════════════════════
def slide_thankyou(prs):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_NAVY)

    # Decorative circles
    c1 = slide.shapes.add_shape(9, Inches(9.8), Inches(-1.5), Inches(6), Inches(6))
    c1.fill.solid(); c1.fill.fore_color.rgb = ACCENT2; c1.line.fill.background()
    c2 = slide.shapes.add_shape(9, Inches(-1.5), Inches(4.5), Inches(4), Inches(4))
    c2.fill.solid(); c2.fill.fore_color.rgb = ACCENT2; c2.line.fill.background()

    # Divider
    add_rect(slide, 2.5, 3.6, 8.33, 0.06, fill_color=ACCENT)

    add_textbox(slide, "Thank You!", 0, 1.38, 13.33, 1.1,
                font_size=52, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, font_name="Calibri Light")
    add_textbox(slide,
                "Multi-Attack Honeypot Network Security System",
                0, 2.55, 13.33, 0.6,
                font_size=18, italic=True, color=ACCENT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "Open for Questions & Discussion",
                0, 3.8, 13.33, 0.6,
                font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Team cards
    for i, (name, usn) in enumerate([
        ("Aditya Kumar",   "1RVU23CSE029"),
        ("Aditi Gopinath", "1RVU23CSE026"),
    ]):
        x = 2.5 + i * 4.8
        add_rect(slide, x, 4.62, 4.3, 1.55, fill_color=ACCENT2)
        add_rect(slide, x, 4.62, 4.3, 0.06, fill_color=ACCENT)
        add_textbox(slide, name, x+0.18, 4.75, 3.9, 0.5,
                    font_size=18, bold=True, color=WHITE)
        add_textbox(slide, usn,  x+0.18, 5.22, 3.9, 0.42,
                    font_size=13, color=ACCENT)
        add_textbox(slide, "B.Tech CSE  •  RV University",
                    x+0.18, 5.62, 3.9, 0.38,
                    font_size=11, italic=True, color=MID_GRAY)

    add_textbox(slide, "github.com/Vinay-R-S/HoneyPot-Network-Security",
                0, 6.55, 13.33, 0.42,
                font_size=12, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════════════════════════
slide_title(prs)
slide_agenda(prs)
slide_intro(prs)
slide_architecture(prs)
slide_attacks(prs)
slide_ml(prs)
slide_simulation(prs)
slide_dashboard(prs)
slide_stack(prs)
slide_results(prs)
slide_conclusion(prs)
slide_thankyou(prs)

out = "HoneyPot_NS_Presentation.pptx"
prs.save(out)
print(f"[OK] Saved -> {out}  ({len(prs.slides)} slides)")
